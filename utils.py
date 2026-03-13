import copy
from pptx.oxml.ns import qn


def _copy_slide(src_prs, src_idx: int, dst_prs, insert_idx: int = None):
    """
    Deep-copy a slide from src_prs into dst_prs, preserving images and background.

    Args:
        src_prs:    Source Presentation object.
        src_idx:    Zero-based index of the slide to copy.
        dst_prs:    Destination Presentation object (can be same as src_prs).
        insert_idx: Position to insert the new slide (0-based). Appends if None.

    Returns:
        The newly created slide object.
    """
    src_slide    = src_prs.slides[src_idx]
    blank_layout = dst_prs.slide_layouts[6]
    new_slide    = dst_prs.slides.add_slide(blank_layout)

    # Remap image relationships
    rId_map = {}
    for rId, rel in src_slide.part.rels.items():
        if "image" in rel.reltype:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

    # Deep-copy slide XML
    slide_xml = copy.deepcopy(src_slide._element)
    if rId_map:
        from lxml import etree
        xml_str = etree.tostring(slide_xml, encoding="unicode")
        for old_id, new_id in rId_map.items():
            xml_str = xml_str.replace(f'r:embed="{old_id}"', f'r:embed="{new_id}"')
            xml_str = xml_str.replace(f'r:id="{old_id}"',    f'r:id="{new_id}"')
        slide_xml = etree.fromstring(xml_str)

    # Copy shape tree
    new_slide.shapes._spTree.clear()
    src_spTree = slide_xml.find(".//" + qn("p:spTree"))
    if src_spTree is not None:
        for el in src_spTree:
            new_slide.shapes._spTree.append(copy.deepcopy(el))

    # Copy background
    src_cSld = slide_xml.find(qn("p:cSld"))
    dst_cSld = new_slide._element.find(qn("p:cSld"))
    if src_cSld is not None and dst_cSld is not None:
        bg = src_cSld.find(qn("p:bg"))
        if bg is not None:
            dst_cSld.insert(0, copy.deepcopy(bg))

    # Reposition slide to insert_idx
    if insert_idx is not None:
        xml_slides = dst_prs.slides._sldIdLst
        inserted   = xml_slides[-1]
        xml_slides.remove(inserted)
        xml_slides.insert(insert_idx, inserted)

    return new_slide