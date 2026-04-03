import os
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from config import BLUEPRINT_IDX, TEMPLATE_FILE, SMM_METRICS, TV_METRICS_BY_CHANNEL
from utils import _copy_slide
from ppt_helpers import (
    add_centered_textbox,
    build_channel_kpi_cards,
    build_kpi_cards_grid,
    build_smm_channel_cards,
    build_youtube_rubric_cards,
    build_totals_slide,
    is_valid_value,
    style_card,
    style_text_frame,
    START_TOP, START_LEFT, CARD_W, CARD_H, CARD_GAP_X, HEADER_COLOR,
)
from pptx.util import Pt

YELLOW = RGBColor(249, 214, 5)
WHITE  = RGBColor(255, 255, 255)


def _df_has_data(df) -> bool:
    """Check if a dataframe has at least one valid non-zero non-null value."""
    if df is None or df.empty:
        return False
    if "Value" in df.columns:
        return df["Value"].apply(is_valid_value).any()
    # Channel / Rubric pivot — check all non-index columns
    value_cols = [c for c in df.columns if c not in ("Channel", "Rubric", "Metric")]
    return any(df[col].apply(is_valid_value).any() for col in value_cols)


def build_presentation(partner: str, country: str, month: str, data: dict, output_path: str):
    """
    Build the full presentation for one partner and save it.
    Skips platform sections entirely if no data exists for that platform.
    """
    # ── Replace if exists
    if os.path.exists(output_path):
        os.remove(output_path)

    prs     = Presentation(TEMPLATE_FILE)
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    offset = [0]

    def copy():
        offset[0] += 1
        return _copy_slide(prs, src_idx=BLUEPRINT_IDX, dst_prs=prs, insert_idx=BLUEPRINT_IDX + offset[0])

    def title_slide(text: str):
        slide = copy()
        add_centered_textbox(slide, text=text, font_size=54, font_color=YELLOW,
                             top=SLIDE_H * 0.38, height=Inches(1.0), slide_w=SLIDE_W)

    # ──────────────────────────────────────────────
    # SLIDE 1 — Partner / Country / Month (always created)
    # ──────────────────────────────────────────────
    slide1 = copy()
    add_centered_textbox(slide1, text=partner, font_size=40, font_color=YELLOW,
                         top=SLIDE_H * 0.30, height=Inches(0.8), slide_w=SLIDE_W)
    add_centered_textbox(slide1, text=f"{country}  |  {month}", font_size=24, font_color=WHITE,
                         top=SLIDE_H * 0.42, height=Inches(0.6), slide_w=SLIDE_W, bold=False)

    # ──────────────────────────────────────────────
    # TV section — check pivoted dataframes directly
    # ──────────────────────────────────────────────
    tv_has_data = _df_has_data(data["tv_channel"]) or _df_has_data(data["tv_other"])
    if tv_has_data:
        title_slide("TV REPORT")
        tv_slide = copy()
        build_channel_kpi_cards(tv_slide, data["tv_channel"], TV_METRICS_BY_CHANNEL)

        right_start = START_LEFT + len(TV_METRICS_BY_CHANNEL) * (CARD_W + CARD_GAP_X) + Inches(0.3)
        visible_k   = 0
        for _, row_other in data["tv_other"].iterrows():
            value = row_other["Value"]
            if not is_valid_value(value):
                continue
            top  = START_TOP + Inches(0.5) + visible_k * (CARD_H + Inches(0.1))
            card = tv_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_start, top, CARD_W, CARD_H)
            style_card(card)
            card.text_frame.text = f"{row_other['Metric']}\n{value}"
            style_text_frame(card.text_frame, font_size=12)
            visible_k += 1

    # ──────────────────────────────────────────────
    # OTT section — check pivoted dataframe directly
    # ──────────────────────────────────────────────
    if _df_has_data(data["ott"]):
        print("ott data exists")
        title_slide("OTT REPORT")
        build_kpi_cards_grid(copy(), data["ott"], SLIDE_W, SLIDE_H)

    # ──────────────────────────────────────────────
    # SMM section — check pivoted dataframes directly
    # ──────────────────────────────────────────────
    smm_has_data = _df_has_data(data["smm_channel"])
    yt_has_data  = _df_has_data(data["yt_rubric"])

    if smm_has_data or yt_has_data:
        title_slide("SOCIAL MEDIA REPORT")

        if smm_has_data:
            build_smm_channel_cards(copy(), data["smm_channel"], SMM_METRICS)

        if yt_has_data:
            build_youtube_rubric_cards(copy(), data["yt_rubric"], SMM_METRICS)

    # ──────────────────────────────────────────────
    # Totals — only if any platform has data
    # ──────────────────────────────────────────────
    if tv_has_data or _df_has_data(data["ott"]) or smm_has_data or yt_has_data:
        build_totals_slide(copy(), data["all_totals"], SLIDE_W)

    # ──────────────────────────────────────────────
    # Delete blueprint and save
    # ──────────────────────────────────────────────
    del prs.slides._sldIdLst[BLUEPRINT_IDX]
    prs.save(output_path)