import math
import pandas as pd
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from config import (
    CARD_W, CARD_H, CARD_GAP_X, CARD_GAP_Y,
    START_TOP, START_LEFT, RUBRIC_LABEL_W, RUBRIC_GAP,
    BORDER_COLOR, TEXT_COLOR, HEADER_COLOR, CARD_BG_COLOR,
)


# ──────────────────────────────────────────────
# Value validation
# ──────────────────────────────────────────────
def is_valid_value(value) -> bool:
    """Return True if value exists, is non-zero, and is not NaN."""
    if value is None:
        return False
    if pd.isna(value):
        return False
    if value == 0:
        return False
    return True


# ──────────────────────────────────────────────
# Card styling
# ──────────────────────────────────────────────
def style_card(card, bg_color=CARD_BG_COLOR):
    """Apply consistent card styling (fill + border)."""
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.5)


def style_text_frame(text_frame, font_size: int):
    """Apply consistent font styling to all paragraphs in a text frame."""
    for para in text_frame.paragraphs:
        para.font.size = Pt(font_size)
        para.font.bold = True
        para.font.color.rgb = TEXT_COLOR
        para.alignment = PP_ALIGN.CENTER


# ──────────────────────────────────────────────
# Textbox helpers
# ──────────────────────────────────────────────
def add_centered_textbox(slide, text: str, font_size: int, font_color: RGBColor,
                         top, height, slide_w, bold: bool = True):
    """Add a full-width centered textbox at a given vertical position."""
    txBox = slide.shapes.add_textbox(Inches(0), top, slide_w, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = font_color


def add_header_label(slide, text: str, left, top, width=Inches(9)):
    """Add a channel/section header label."""
    box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HEADER_COLOR


# ──────────────────────────────────────────────
# Card builders
# ──────────────────────────────────────────────
def add_card(slide, left, top, width, metric: str, value, font_size: int):
    """Add a single KPI card to a slide."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, CARD_H)
    style_card(card)
    formatted_value = f"{int(value):,}" if is_valid_value(value) else "0"
    card.text_frame.text = f"{metric}\n{formatted_value}"
    style_text_frame(card.text_frame, font_size)


def build_channel_kpi_cards(slide, df_pivot, metrics_list: list):
    """
    Build channel header + KPI cards layout.
    Each channel gets a header label and a row of metric cards below it.
    Used for TV Summary and SMM Summary.
    """
    for i, row in df_pivot.iterrows():
        header_top = START_TOP + i * (CARD_H + CARD_GAP_Y)
        add_header_label(slide, text=row["Channel"], left=START_LEFT, top=header_top)

        visible_j = 0
        for metric in metrics_list:
            value = row.get(metric, 0)
            if not is_valid_value(value):
                continue
            left = START_LEFT + visible_j * (CARD_W + CARD_GAP_X)
            top  = header_top + Inches(0.5)
            add_card(slide, left, top, CARD_W, metric, value, font_size=14)
            visible_j += 1

def build_smm_channel_cards(slide, df_pivot, metrics_list: list):
    """
    Build SMM channel rubric rows.
    2 channels: stacked vertically (full width).
    3-4 channels: 2x2 grid, slide split into 4 equal quadrants.
    """
    channels = list(df_pivot.groupby("Channel", sort=False))
    n_channels = len(channels)

    # Slide dimensions (standard 13.333 x 7.5 inches)
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    HALF_W = SLIDE_W / 2
    HALF_H = SLIDE_H / 2

    PADDING = Inches(0.3)

    if n_channels <= 2:
        _draw_channel_column(slide, channels, metrics_list,
                             left_offset=START_LEFT, top_offset=START_TOP,
                             rubric_w=RUBRIC_LABEL_W)
    else:
        # 2x2 grid positions: (col, row)
        positions = [
            (PADDING, PADDING),                          # top-left
            (HALF_W + PADDING, PADDING),                 # top-right
            (PADDING, HALF_H + PADDING),                 # bottom-left
            (HALF_W + PADDING, HALF_H + PADDING),        # bottom-right
        ]

        for idx, (channel_name, group) in enumerate(channels):
            if idx >= 4:
                break
            left, top = positions[idx]
            _draw_channel_quadrant(slide, channel_name, group, metrics_list,
                                   left_offset=left, top_offset=top,
                                   rubric_w=Inches(0.8))


def _draw_channel_quadrant(slide, channel_name, group, metrics_list,
                           left_offset, top_offset, rubric_w=None):
    """
    Draw a single channel block in its quadrant.
    """
    if rubric_w is None:
        rubric_w = RUBRIC_LABEL_W

    current_top = top_offset

    # Channel header
    box = slide.shapes.add_textbox(left_offset, current_top, Inches(5), Inches(0.5))
    tf = box.text_frame
    tf.text = channel_name
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HEADER_COLOR

    current_top += Inches(0.5)

    for _, row in group.iterrows():
        # Content label
        rubric_box = slide.shapes.add_textbox(left_offset, current_top, rubric_w, CARD_H)
        tf = rubric_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = row["Content"]
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = TEXT_COLOR

        cards_start_left = left_offset + rubric_w + RUBRIC_GAP
        visible_j = 0
        for metric in metrics_list:
            value = row.get(metric, 0)
            if not is_valid_value(value):
                continue
            card_left = cards_start_left + visible_j * (CARD_W + CARD_GAP_X)
            add_card(slide, card_left, current_top, CARD_W, metric, value, font_size=12)
            visible_j += 1

        current_top += CARD_H + Inches(0.15)


def _draw_channel_column(slide, channels, metrics_list, left_offset, top_offset, rubric_w=None):
    """
    Draw channels stacked vertically (for 2 or fewer channels).
    """
    if rubric_w is None:
        rubric_w = RUBRIC_LABEL_W

    current_top = top_offset

    for channel_name, group in channels:
        _draw_channel_quadrant(slide, channel_name, group, metrics_list,
                               left_offset=left_offset, top_offset=current_top,
                               rubric_w=rubric_w)

        row_count = len(group)
        current_top += Inches(0.5) + row_count * (CARD_H + Inches(0.15)) + Inches(0.2)


def build_kpi_cards_grid(slide, df_pivot, slide_w, slide_h):
    """
    Build KPI cards in a centered 2-row grid layout.
    Row split: 10→5+5 | 9→5+4 | 8→4+4 | 7→4+3 ...
    Used for OTT Summary.
    """
    valid_rows = [
        (row["Metric"], row["Value"])
        for _, row in df_pivot.iterrows()
        if is_valid_value(row["Value"])
    ]
    if not valid_rows:
        return

    total_cards = len(valid_rows)
    row1_count  = math.ceil(total_cards / 2)
    rows        = [valid_rows[:row1_count], valid_rows[row1_count:]]
    row_tops    = [slide_h * 0.30, slide_h * 0.30 + CARD_H + Inches(0.2)]

    for row_idx, row_cards in enumerate(rows):
        n = len(row_cards)
        if n == 0:
            continue
        total_row_width = n * CARD_W + (n - 1) * CARD_GAP_X
        row_start_left  = (slide_w - total_row_width) / 2
        top             = row_tops[row_idx]
        for col_idx, (metric, value) in enumerate(row_cards):
            left = row_start_left + col_idx * (CARD_W + CARD_GAP_X)
            add_card(slide, left, top, CARD_W, metric, value, font_size=12)


def build_youtube_rubric_cards(slide, df_pivot, metrics_list: list):
    """
    Build YouTube rubric rows: rubric label on left, KPI cards to the right.
    """
    # Youtube header
    box = slide.shapes.add_textbox(START_LEFT, START_TOP, Inches(9), Inches(0.5))
    tf = box.text_frame
    tf.text = "Youtube"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HEADER_COLOR

    for i, row in df_pivot.iterrows():
        row_top = START_TOP + Inches(0.6) + i * (CARD_H + Inches(0.15))

        # Rubric label — vertically centered on the left
        rubric_box = slide.shapes.add_textbox(START_LEFT, row_top, RUBRIC_LABEL_W, CARD_H)
        tf = rubric_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = row["Content"]
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = TEXT_COLOR

        cards_start_left = START_LEFT + RUBRIC_LABEL_W + RUBRIC_GAP
        visible_j = 0
        for metric in metrics_list:
            value = row.get(metric, 0)
            if not is_valid_value(value):
                continue
            left = cards_start_left + visible_j * (CARD_W + CARD_GAP_X)
            add_card(slide, left, row_top, CARD_W, metric, value, font_size=12)
            visible_j += 1


def build_totals_slide(slide, df_all_totals, slide_w):
    """
    Build Totals slide with wrapped rows per platform.
    Platform label on the left above the KPI cards block.
    """
    BLOCK_GAP    = Inches(0.3)
    HEADER_H     = Inches(0.4)
    ROW_GAP      = Inches(0.1)
    current_top  = START_TOP

    # Max cards per row based on slide width
    available_width = slide_w - (START_LEFT * 2)
    max_per_row     = int(available_width // (CARD_W + CARD_GAP_X))

    for platform in ["TV", "OTT", "SMM"]:
        df_platform = df_all_totals[df_all_totals["Platform"] == platform].reset_index(drop=True)

        valid_rows = [
            (row["Metric"], row["Value"])
            for _, row in df_platform.iterrows()
            if is_valid_value(row["Value"])
        ]
        if not valid_rows:
            continue

        # Platform label — left aligned, above the cards
        box = slide.shapes.add_textbox(START_LEFT, current_top, Inches(3), HEADER_H)
        tf  = box.text_frame
        tf.word_wrap = False
        tf.text = platform
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = HEADER_COLOR
        current_top += HEADER_H

        # Split into rows of max_per_row and render each row
        chunks = [valid_rows[i:i + max_per_row] for i in range(0, len(valid_rows), max_per_row)]
        for chunk in chunks:
            for col_idx, (metric, value) in enumerate(chunk):
                left = START_LEFT + col_idx * (CARD_W + CARD_GAP_X)
                add_card(slide, left, current_top, CARD_W, metric, value, font_size=11)
            current_top += CARD_H + ROW_GAP

        current_top += BLOCK_GAP